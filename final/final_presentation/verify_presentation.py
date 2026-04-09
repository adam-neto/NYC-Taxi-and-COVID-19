from __future__ import annotations

import hashlib
import zipfile
from pathlib import Path

from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent
PPTX_PATH = BASE_DIR / "group10_final_presentation_first_draft.pptx"

EXPECTED_TITLES = {
    1: "COVID-19 and Behavioral Shifts in NYC Taxi Travel",
    2: "Why This Problem Matters",
    3: "Research Questions and Contributions",
    4: "Dataset Statistics and Period Design",
    5: "Shared Workflow and Methodology",
    6: "RQ1: Tipping Behavior",
    7: "RQ2: Airport Recovery Descriptively",
    8: "RQ2: Regression Evidence",
    9: "RQ3: Cashless Shift Descriptively",
    10: "RQ3: Model-Based Evaluation",
    11: "Discussion, What Did Not Work, and Limitations",
    12: "Conclusion",
    13: "Backup: RQ2 Coefficient Table",
    14: "Backup: RQ3 Full Metrics Table",
    15: "Backup: Replication Workflow",
}

EXPECTED_SNIPPETS = {
    6: ["22.57%", "23.34%", "23.35%", "22.19%", "+0.764 pp", "+0.775 pp", "-0.388 pp"],
    7: ["JFK: 0.704", "LaGuardia: 0.577"],
    8: ["+0.0108, p = 0.0031", "+0.0125, p < 1e-9"],
    9: ["72.59%", "73.96%", "78.46%", "82.34%", "+1.38 pp", "+5.87 pp", "+9.75 pp"],
    10: ["0.500", "0.581", "0.638", "+16.25%", "+27.54%", "513,206", "56,105"],
    13: ["+0.0005, p = 0.9175", "+0.0108, p = 0.0031", "+0.0125, p < 1e-9"],
}

FIGURE_PATHS = [
    BASE_DIR.parent.parent / "RQ1/figures/monthly_tipping_trend.png",
    BASE_DIR.parent.parent / "RQ1/figures/period_tipping_change.png",
    BASE_DIR.parent.parent / "RQ2/figures/monthly_airport_trip_counts.png",
    BASE_DIR.parent.parent / "RQ2/figures/monthly_airport_trip_share.png",
    BASE_DIR.parent.parent / "RQ2/figures/airport_regression_interactions.png",
    BASE_DIR.parent.parent / "RQ3/figures/monthly_cashless_share.png",
    BASE_DIR.parent.parent / "RQ3/figures/period_payment_mix.png",
    BASE_DIR.parent.parent / "RQ3/figures/cashless_model_metric_comparison.png",
    BASE_DIR.parent.parent / "RQ3/figures/cashless_xgboost_feature_importance.png",
]


def md5_bytes(payload: bytes) -> str:
    return hashlib.md5(payload).hexdigest()


def verify_titles(prs: Presentation) -> list[str]:
    errors: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        title = texts[0].split("\n")[0] if texts else ""
        expected = EXPECTED_TITLES.get(idx)
        if title != expected:
            errors.append(f"Slide {idx} title mismatch: expected '{expected}', found '{title}'")
    return errors


def verify_snippets(prs: Presentation) -> list[str]:
    errors: list[str] = []
    for idx, snippets in EXPECTED_SNIPPETS.items():
        slide = prs.slides[idx - 1]
        text_blob = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        missing = [snippet for snippet in snippets if snippet not in text_blob]
        if missing:
            errors.append(f"Slide {idx} missing snippets: {missing}")
    return errors


def verify_images() -> list[str]:
    errors: list[str] = []
    source_hashes = {path.name: md5_bytes(path.read_bytes()) for path in FIGURE_PATHS}
    media_hashes: list[str] = []
    with zipfile.ZipFile(PPTX_PATH, "r") as zf:
        for name in zf.namelist():
            if name.startswith("ppt/media/"):
                media_hashes.append(md5_bytes(zf.read(name)))
    for name, digest in source_hashes.items():
        if digest not in media_hashes:
            errors.append(f"Embedded image missing or changed: {name}")
    return errors


def main() -> int:
    prs = Presentation(str(PPTX_PATH))
    errors: list[str] = []
    errors.extend(verify_titles(prs))
    errors.extend(verify_snippets(prs))
    errors.extend(verify_images())

    if errors:
        print("Presentation verification failed:")
        for err in errors:
            print(f"- {err}")
        return 1

    print("Presentation verification passed.")
    print(f"Verified file: {PPTX_PATH}")
    print(f"Checked {len(prs.slides)} slides and {len(FIGURE_PATHS)} source figures.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
