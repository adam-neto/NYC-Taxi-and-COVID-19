from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PROJECT_DIR = BASE_DIR.parent.parent
OUTPUT_PATH = BASE_DIR / "group10_final_presentation_first_draft.pptx"


FIGURES = {
    "rq1_monthly": PROJECT_DIR / "RQ1/figures/monthly_tipping_trend.png",
    "rq1_period": PROJECT_DIR / "RQ1/figures/period_tipping_change.png",
    "rq2_counts": PROJECT_DIR / "RQ2/figures/monthly_airport_trip_counts.png",
    "rq2_share": PROJECT_DIR / "RQ2/figures/monthly_airport_trip_share.png",
    "rq2_reg": PROJECT_DIR / "RQ2/figures/airport_regression_interactions.png",
    "rq3_monthly": PROJECT_DIR / "RQ3/figures/monthly_cashless_share.png",
    "rq3_mix": PROJECT_DIR / "RQ3/figures/period_payment_mix.png",
    "rq3_metrics": PROJECT_DIR / "RQ3/figures/cashless_model_metric_comparison.png",
    "rq3_xgb": PROJECT_DIR / "RQ3/figures/cashless_xgboost_feature_importance.png",
}


BG = RGBColor(248, 249, 251)
NAVY = RGBColor(20, 33, 61)
BLUE = RGBColor(48, 94, 180)
RED = RGBColor(178, 42, 43)
GOLD = RGBColor(214, 161, 46)
TEXT = RGBColor(45, 52, 54)
MUTED = RGBColor(110, 117, 126)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(231, 239, 252)
LIGHT_GOLD = RGBColor(251, 245, 226)
LIGHT_RED = RGBColor(252, 237, 237)


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number(slide, number: int):
    box = slide.shapes.add_textbox(Inches(12.65), Inches(7.0), Inches(0.45), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = MUTED
    run.font.name = "Aptos"


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9.7), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.size = Pt(27)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Aptos Display"
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.93), Inches(11.8), Inches(0.35))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        srun = sp.runs[0]
        srun.font.size = Pt(12)
        srun.font.color.rgb = MUTED
        srun.font.name = "Aptos"


def add_bullets(slide, x, y, w, h, bullets, font_size=20, color=TEXT, level0_bold=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.space_after = Pt(7)
        p.bullet = True
        run = p.runs[0]
        run.font.size = Pt(font_size - 1 if level else font_size)
        run.font.color.rgb = color
        run.font.name = "Aptos"
        if level0_bold and level == 0:
            run.font.bold = True
    return box


def add_callout(slide, x, y, w, h, title, lines, fill_color=LIGHT_BLUE, title_color=NAVY):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = title
    p.space_after = Pt(6)
    run = p.runs[0]
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = "Aptos"
    for line in lines:
        q = tf.add_paragraph()
        q.text = line
        q.bullet = True
        q.space_after = Pt(3)
        r = q.runs[0]
        r.font.size = Pt(15)
        r.font.color.rgb = TEXT
        r.font.name = "Aptos"
    return shape


def add_picture(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        raise FileNotFoundError(path)
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_notes(slide, presenter: str, timing: str, script: str):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    p = notes_frame.paragraphs[0]
    p.text = f"Presenter: {presenter}"
    q = notes_frame.add_paragraph()
    q.text = f"Target time: {timing}"
    r = notes_frame.add_paragraph()
    r.text = "Script:"
    s = notes_frame.add_paragraph()
    s.text = script


def make_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(12.0), Inches(1.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "COVID-19 and Behavioral Shifts in NYC Taxi Travel"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Aptos Display"

    p2 = tf.add_paragraph()
    p2.text = "Tipping, Airport Demand, and the Move to Cashless Payments"
    r2 = p2.runs[0]
    r2.font.size = Pt(22)
    r2.font.color.rgb = BLUE
    r2.font.name = "Aptos"

    sub = slide.shapes.add_textbox(Inches(0.72), Inches(3.0), Inches(10.8), Inches(0.7))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Group 10 | Adam Neto, Harsh Kalyani, Chris Wang"
    sr = sp.runs[0]
    sr.font.size = Pt(20)
    sr.font.color.rgb = TEXT
    sr.font.name = "Aptos"

    thesis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(4.0), Inches(11.2), Inches(1.05))
    thesis.fill.solid()
    thesis.fill.fore_color.rgb = LIGHT_BLUE
    thesis.line.color.rgb = LIGHT_BLUE
    ttf = thesis.text_frame
    tp = ttf.paragraphs[0]
    tp.text = "COVID changed NYC yellow taxi behavior, but the recovery was uneven across tipping, airport demand, and payment choice."
    tr = tp.runs[0]
    tr.font.size = Pt(20)
    tr.font.bold = True
    tr.font.color.rgb = NAVY
    tr.font.name = "Aptos"

    footer = slide.shapes.add_textbox(Inches(0.72), Inches(6.45), Inches(4.4), Inches(0.35))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "CISC 351 | Winter 2026"
    fr = fp.runs[0]
    fr.font.size = Pt(14)
    fr.font.color.rgb = MUTED
    fr.font.name = "Aptos"

    add_slide_number(slide, 1)
    add_notes(
        slide,
        "Adam",
        "0:40",
        "We are Group 10, and our project studies how COVID-19 changed New York City yellow taxi behavior before, during, and after the main disruption period. We study three dimensions of that change: tipping behavior, airport-related recovery, and the shift from cash to cashless payments. Our central question is not simply whether taxi demand fell, because that part is obvious. What matters is which behaviors returned to baseline and which remained changed after recovery. Taken together, these three questions let us compare one behavior that mostly reverted, one that recovered unevenly, and one that appears to have shifted more persistently.",
    )


def add_text_slide(prs: Presentation, number: int, title: str, presenter: str, timing: str, bullets, script: str, callout=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, f"Presenter: {presenter} | Target time: {timing}")
    add_bullets(slide, Inches(0.75), Inches(1.45), Inches(7.5), Inches(5.4), bullets, font_size=20)
    if callout:
        add_callout(slide, Inches(8.75), Inches(1.65), Inches(3.8), Inches(3.8), callout["title"], callout["lines"], callout.get("fill", LIGHT_BLUE))
    add_slide_number(slide, number)
    add_notes(slide, presenter, timing, script)
    return slide


def add_two_image_slide(
    prs: Presentation,
    number: int,
    title: str,
    presenter: str,
    timing: str,
    left_img: Path,
    right_img: Path,
    left_caption: str,
    right_caption: str,
    bullets,
    script: str,
    callout=None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, f"Presenter: {presenter} | Target time: {timing}")
    add_bullets(slide, Inches(0.75), Inches(1.35), Inches(4.0), Inches(2.6), bullets, font_size=18)
    if callout:
        add_callout(slide, Inches(0.78), Inches(4.2), Inches(4.0), Inches(2.0), callout["title"], callout["lines"], callout.get("fill", LIGHT_BLUE))
    add_picture(slide, left_img, Inches(5.0), Inches(1.38), w=Inches(3.85))
    add_picture(slide, right_img, Inches(8.95), Inches(1.38), w=Inches(3.85))
    lcap = slide.shapes.add_textbox(Inches(5.05), Inches(4.95), Inches(3.75), Inches(0.6))
    lp = lcap.text_frame.paragraphs[0]
    lp.text = left_caption
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.runs[0]
    lr.font.size = Pt(12)
    lr.font.color.rgb = MUTED
    lr.font.name = "Aptos"
    rcap = slide.shapes.add_textbox(Inches(9.0), Inches(4.95), Inches(3.75), Inches(0.6))
    rp = rcap.text_frame.paragraphs[0]
    rp.text = right_caption
    rp.alignment = PP_ALIGN.CENTER
    rr = rp.runs[0]
    rr.font.size = Pt(12)
    rr.font.color.rgb = MUTED
    rr.font.name = "Aptos"
    add_slide_number(slide, number)
    add_notes(slide, presenter, timing, script)
    return slide


def add_round_robin_discussion_slide(prs: Presentation, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Discussion and Limitations", "Presenters: Adam, Harsh, Chris | Target time: 1:15")

    col_x = [0.65, 4.55, 8.45]
    headers = ["RQ1 Measurement", "RQ2 Outcome", "RQ3 Evaluation"]
    fills = [LIGHT_GOLD, LIGHT_BLUE, LIGHT_RED]
    bodies = [
        ["Tipping is measured on credit-card trips", "Cash tips are not fully observed", "Interpretation is limited to recorded card tips"],
        ["Airport trip share is the primary outcome", "This accounts for swings in total yellow taxi demand", "Robustness checks are also reported"],
        ["ROC-AUC and average precision are emphasized", "The test set is imbalanced toward cashless trips", "Accuracy alone is not enough"],
    ]

    for idx in range(3):
        add_callout(slide, Inches(col_x[idx]), Inches(1.55), Inches(3.55), Inches(3.7), headers[idx], bodies[idx], fills[idx])

    lim = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.65), Inches(11.95), Inches(0.85))
    lim.fill.solid()
    lim.fill.fore_color.rgb = RGBColor(236, 240, 241)
    lim.line.color.rgb = RGBColor(236, 240, 241)
    tf = lim.text_frame
    p = tf.paragraphs[0]
    p.text = "Broader limitations: yellow taxis only, not all NYC travel; administrative data constraints; comparative evidence rather than causal claims."
    r = p.runs[0]
    r.font.size = Pt(17)
    r.font.color.rgb = TEXT
    r.font.name = "Aptos"
    add_slide_number(slide, number)
    add_notes(
        slide,
        "Shared round-robin",
        "1:15",
        "Adam: For RQ2, the main comparative outcome is airport trip share, because total yellow taxi demand also changes a lot across the study window. Looking only at counts would mix airport recovery with the broader recovery of the yellow taxi system.\nHarsh: For RQ1, tipping is measured on credit-card trips, because cash tips are not fully observed in the TLC records. That means the result is informative, but it should be interpreted as recorded card-tip behavior rather than all tipping.\nChris: For RQ3, ROC-AUC and average precision are the most informative model metrics here, because the held-out test set is imbalanced toward cashless trips. A model can look strong on raw accuracy even if it is mostly just tracking the majority class.\nAdam close: Across all three questions, the broader limitation is that this is evidence about yellow taxi behavior, not a causal statement about all travel in New York City. Other factors like fare changes, trip mix, and broader market changes may also contribute to the observed patterns.",
    )


def add_backup_metrics_slide(prs: Presentation, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Backup: RQ3 Full Metrics Table", "Use if asked about evaluation details")
    add_picture(slide, FIGURES["rq3_metrics"], Inches(0.8), Inches(1.45), w=Inches(6.0))
    add_callout(
        slide,
        Inches(7.1),
        Inches(1.6),
        Inches(5.4),
        Inches(3.8),
        "Held-out test set",
        [
            "Train rows: 513,206",
            "Test rows: 56,105",
            "Dummy ROC-AUC: 0.500",
            "Logistic ROC-AUC: 0.581",
            "XGBoost ROC-AUC: 0.638",
        ],
        LIGHT_RED,
    )
    add_slide_number(slide, number)
    add_notes(slide, "Any", "Backup", "Use this slide if the audience asks why ROC-AUC matters more than raw accuracy or wants the exact held-out model metrics.")


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs)

    add_text_slide(
        prs,
        2,
        "Why This Problem Matters",
        "Adam",
        "1:05",
        [
            "COVID changed more than total trip volume",
            "Taxi data reveals behavior inside trips: when people travel, how they tip, how they pay, and what types of trips recover",
            "Different travel behaviors recovered at different rates",
            "Yellow taxi records provide trip-level evidence beyond broad citywide anecdotes",
        ],
        "This problem matters because COVID did not only reduce mobility, it also changed how people traveled, how they paid, and how different types of trips recovered. A simple citywide trip-count graph would show the shock, but it would miss the behavioral detail inside each trip. Yellow taxi data is useful because it gives trip-level information on time, fare, tips, payment type, and pickup and dropoff locations. That means we can study whether tipping changed, whether airport travel came back in the same way across airports, and whether payment habits kept shifting after the city reopened. It also lets us separate the immediate disruption in 2020 from the slower reopening years, which is important if we want to distinguish temporary shocks from more persistent behavioral change.",
        {"title": "Core lens", "lines": ["Trip-level TLC records", "Behavior within trips", "Recovery over 2019-2023"], "fill": LIGHT_GOLD},
    )

    add_text_slide(
        prs,
        3,
        "Research Questions and Contributions",
        "Chris",
        "1:10",
        [
            "RQ1: How did recorded tipping behavior change across the four study periods?",
            "RQ2: Did JFK and LaGuardia recover differently after the COVID shock?",
            "RQ3: Did COVID accelerate the shift from cash to cashless payments, and did that change persist?",
            "Shared contributions: common 4-period design, reproducible DuckDB pipeline, descriptive + inferential + ML evidence",
        ],
        "Our first research question asks how recorded tipping behavior changed over time and whether it eventually returned to the 2019 baseline. Our second asks whether airport-related yellow taxi recovery differed between JFK and LaGuardia, rather than treating airport demand as one combined category. Our third asks whether COVID accelerated the shift from cash to cashless payments and whether that change persisted. What ties these three questions together is a common four-period design, the same TLC data source, and the same reproducible DuckDB pipeline. The questions also give us three different types of evidence: a descriptive tipping analysis, a descriptive plus inferential airport analysis, and a descriptive plus machine-learning payment analysis.",
        {"title": "Four periods", "lines": ["2019 pre-COVID", "2020 COVID", "2021-2022 intermediate", "2023 post-COVID"]},
    )

    add_text_slide(
        prs,
        4,
        "Dataset Statistics and Period Design",
        "Harsh",
        "1:10",
        [
            "NYC TLC yellow taxi trip records from January 2019 through December 2023",
            "60 monthly parquet files plus the taxi zone lookup table",
            "218,118,168 total yellow taxi trips in the study window",
            "Key fields: timestamps, fare, tip, payment type, passenger count, trip distance, pickup and dropoff zones",
            "Four periods: pre-COVID, COVID, intermediate recovery, post-COVID",
        ],
        "Our dataset is the New York City TLC yellow taxi trip record data from January 2019 through December 2023, plus the taxi zone lookup table. In total, the study window contains 60 monthly parquet files and about 218.1 million yellow taxi trips. The variables we rely on most are pickup and dropoff time, pickup and dropoff zone, fare amount, tip amount, total amount, trip distance, passenger count, and payment type. We organize the analysis into four periods: pre-COVID in 2019, the disruption year in 2020, an intermediate recovery period in 2021 and 2022, and a post-COVID period in 2023. Including 2021 and 2022 is important because recovery was gradual and uneven. If we jumped directly from 2020 to 2023, we would miss the transition period where some behaviors started to rebound while others were still far from baseline.",
        {"title": "Why 2021-2022 matters", "lines": ["Captures uneven reopening", "Prevents a misleading jump from 2020 straight to 2023"], "fill": LIGHT_GOLD},
    )

    add_text_slide(
        prs,
        5,
        "Shared Workflow and Methodology",
        "Harsh",
        "1:10",
        [
            "Local TLC parquet files are queried through a shared DuckDB layer",
            "`query_taxi_duckdb.py` assigns the common period labels and builds reproducible summaries",
            "RQ1 filters to valid credit-card trips",
            "RQ2 identifies airport-linked trips and estimates a monthly panel regression",
            "RQ3 groups payment types and adds a held-out ML evaluation",
        ],
        "All three research questions use the same DuckDB-based query workflow on the local parquet files instead of loading the entire dataset into one giant in-memory table. The shared script assigns each month to one of the four study periods and produces a consistent base for the downstream analyses. That consistency matters because we want all three questions to be comparable, not built on different period definitions or different preprocessing choices. From there, each research question applies its own filtering and analysis logic. RQ1 focuses on valid credit-card trips because cash tips are not reliably observed. RQ2 identifies airport-related trips using JFK and LaGuardia zone definitions and then builds both descriptive summaries and a monthly panel regression. RQ3 groups payment types into cashless, cash, and ambiguous categories and then adds a held-out machine learning evaluation.",
        {"title": "Pipeline", "lines": ["Preprocess", "Question-specific filtering", "Summaries and models", "Figure generation"]},
    )

    add_two_image_slide(
        prs,
        6,
        "RQ1: Tipping Behavior",
        "Harsh",
        "1:20",
        FIGURES["rq1_monthly"],
        FIGURES["rq1_period"],
        "Monthly average tip percentage on credit-card trips",
        "Change in tip percentage versus the 2019 baseline",
        [
            "Credit-card trips only, because cash tips are not fully observed",
            "2019 baseline: 22.57%",
            "COVID: 23.34% | Intermediate: 23.35% | Post-COVID: 22.19%",
        ],
        "For RQ1, we measure tipping using valid credit-card trips only, because cash tips are not reliably captured in the TLC records. Even with that restriction, the sample is still very large: about 60.6 million valid credit-card trips in 2019, 17.5 million in 2020, 52.6 million in the intermediate period, and 29.9 million in 2023. The baseline average tip percentage in 2019 is 22.57 percent. That rises to 23.34 percent in 2020 and stays at 23.35 percent in the intermediate recovery period, before dropping to 22.19 percent in 2023. Relative to the baseline, that is a positive shift of about 0.76 to 0.78 percentage points during the disruption and intermediate recovery, followed by a decline to about 0.39 percentage points below baseline in 2023. So the overall pattern is not a steady upward trend. It looks more like a temporary increase in recorded card-tip share that mostly fades out by the post-COVID period.",
        {"title": "Change vs. baseline", "lines": ["COVID: +0.764 pp", "Intermediate: +0.775 pp", "Post-COVID: -0.388 pp"], "fill": LIGHT_GOLD},
    )

    add_two_image_slide(
        prs,
        7,
        "RQ2: Airport Recovery Descriptively",
        "Adam",
        "1:30",
        FIGURES["rq2_counts"],
        FIGURES["rq2_share"],
        "Monthly airport-related yellow taxi trip counts",
        "Airport trips as a share of all yellow taxi trips",
        [
            "Comparison: JFK and LaGuardia airport-related recovery",
            "Both airports collapsed in 2020",
            "JFK recovered more strongly in counts and system share",
        ],
        "RQ2 compares airport-related yellow taxi recovery at JFK and LaGuardia. Descriptively, both airports collapse in 2020, but JFK recovers more strongly afterward, both in absolute monthly volume and as a share of the yellow taxi system. Using average monthly trip counts to keep the comparison fair across periods of different lengths, JFK falls from about 281,827 airport-related trips per month in 2019 to about 62,362 in 2020. LaGuardia falls from about 250,675 to about 46,398. During the intermediate recovery period, JFK climbs back to roughly 148,752 trips per month, while LaGuardia recovers to about 101,457. In 2023, JFK reaches about 198,405 trips per month, or 70.4 percent of its 2019 baseline, while LaGuardia reaches about 144,561 trips per month, or 57.7 percent of baseline. The share view tells the same story: JFK rises from about 4.03 percent of all yellow taxi trips before COVID to 6.23 percent in 2023, while LaGuardia rises from 3.56 percent to 4.52 percent. So recovery is not only incomplete, it is also uneven.",
        {"title": "2023 recovery index", "lines": ["JFK: 0.704", "LaGuardia: 0.577"], "fill": LIGHT_BLUE},
    )

    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide8)
    add_title(slide8, "RQ2: Regression Evidence", "Presenters: Adam then Chris | Target time: 1:20")
    add_picture(slide8, FIGURES["rq2_reg"], Inches(6.0), Inches(1.35), w=Inches(6.5))
    add_bullets(
        slide8,
        Inches(0.75),
        Inches(1.4),
        Inches(4.7),
        Inches(2.7),
        [
            "Monthly airport panel regression",
            "LaGuardia and 2019 baseline are the reference categories",
            "Key evidence: JFK-by-period interaction terms",
            "Positive interactions mean stronger JFK recovery relative to LaGuardia",
        ],
        font_size=18,
    )
    add_callout(slide8, Inches(0.78), Inches(4.25), Inches(4.7), Inches(1.75), "Key coefficients", ["JFK x intermediate: +0.0108, p = 0.0031", "JFK x post-COVID: +0.0125, p < 1e-9"], LIGHT_BLUE)
    add_slide_number(slide8, 8)
    add_notes(
        slide8,
        "Adam then Chris",
        "1:20",
        "Adam: To test whether that descriptive gap is statistically meaningful, we estimate a monthly airport panel regression. The outcome is airport trip share, LaGuardia and 2019 are the reference categories, and we include month-of-year fixed effects so that seasonal variation is not confused with recovery. The key coefficients are the JFK-by-period interaction terms, because they tell us whether JFK changes differently from LaGuardia in each later phase.\nChris: The COVID interaction is essentially zero, which means the initial collapse does not look meaningfully different across the two airports once we compare them relative to baseline. The divergence appears later. The intermediate interaction is about positive 0.0108 with a p-value of 0.0031, and the post-COVID interaction is about positive 0.0125 with a p-value below one times ten to the negative ninth. In substantive terms, JFK's airport-trip share rises by roughly 1.1 to 1.3 percentage points more than LaGuardia's in the recovery phases. The pickup-only robustness check shows the same sign and significance pattern, which strengthens the claim that the uneven recovery is not just an artifact of one airport definition.",
    )

    add_two_image_slide(
        prs,
        9,
        "RQ3: Cashless Shift Descriptively",
        "Chris",
        "1:20",
        FIGURES["rq3_monthly"],
        FIGURES["rq3_mix"],
        "Monthly cashless share",
        "Payment mix by project period",
        [
            "Cashless payment was already common before COVID",
            "Cashless usage continues to rise across the study window",
            "Known-payment cashless share: 72.59% -> 73.96% -> 78.46% -> 82.34%",
        ],
        "For RQ3, cashless payment is already common before COVID, but it becomes even more prevalent over the study window. Using cashless share among known payment trips, the baseline is 72.59 percent in 2019. That rises to 73.96 percent in 2020, 78.46 percent in the intermediate recovery period, and 82.34 percent in 2023. Relative to the baseline, that is an increase of about 1.38 percentage points during COVID, 5.87 points in the intermediate period, and 9.75 points by the post-COVID period. So unlike the tipping result, this pattern does not move back toward baseline. It keeps strengthening. We treat known-payment share as the primary measure because ambiguous payment codes vary over time, and we do not want those administrative coding shifts to blur the cash-versus-cashless comparison. The descriptive takeaway is that the move away from cash persists well beyond the initial disruption.",
        {"title": "Change vs. baseline", "lines": ["COVID: +1.38 pp", "Intermediate: +5.87 pp", "Post-COVID: +9.75 pp"], "fill": LIGHT_RED},
    )

    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide10)
    add_title(slide10, "RQ3: Model-Based Evaluation", "Presenters: Harsh then Chris | Target time: 1:25")
    add_picture(slide10, FIGURES["rq3_metrics"], Inches(5.8), Inches(1.35), w=Inches(3.2))
    add_picture(slide10, FIGURES["rq3_xgb"], Inches(9.15), Inches(1.35), w=Inches(3.45))
    add_bullets(
        slide10,
        Inches(0.75),
        Inches(1.4),
        Inches(4.5),
        Inches(2.5),
        [
            "Known-payment trips only",
            "Train before July 2023, test from July 2023 onward",
            "Train rows: 513,206 | Test rows: 56,105",
            "Compare dummy baseline, logistic regression, and XGBoost",
        ],
        font_size=18,
    )
    add_callout(slide10, Inches(0.78), Inches(4.15), Inches(4.5), Inches(2.2), "ROC-AUC comparison", ["Dummy: 0.500", "Logistic: 0.581 (+16.25%)", "XGBoost: 0.638 (+27.54%)"], LIGHT_RED)
    add_slide_number(slide10, 10)
    add_notes(
        slide10,
        "Harsh then Chris",
        "1:25",
        "Harsh: For RQ3, we complement the descriptive analysis with a held-out prediction task. The model predicts whether a known-payment trip is cashless or cash. We train on trips before July 2023 and test on trips from July 2023 onward, which gives us 513,206 training rows and 56,105 held-out test rows. We use a time-based split rather than a random split so the evaluation reflects a more realistic future-prediction setting and reduces temporal leakage. Because the test set is heavily tilted toward cashless trips, accuracy by itself can be misleading.\nChris: That is why ROC-AUC and average precision matter more here. The dummy most-frequent baseline has ROC-AUC of 0.500. Logistic regression improves that to 0.581, and XGBoost improves further to 0.638. Those are improvements of about 16.25 percent and 27.54 percent over the dummy baseline. Average precision also rises from 0.825 for the dummy model to 0.859 for logistic regression and 0.885 for XGBoost. The feature-importance results show that period indicators, borough variables, and fare amount remain among the strongest features, which supports the descriptive conclusion that payment behavior changes across recovery phases rather than only across trip geography.",
    )

    add_round_robin_discussion_slide(prs, 11)

    add_text_slide(
        prs,
        12,
        "Conclusion",
        "Chris",
        "0:40",
        [
            "COVID changed multiple parts of yellow taxi behavior, but those parts did not recover in the same way",
            "Tipping mostly moved back toward baseline by 2023",
            "Airport recovery stayed uneven, with stronger JFK recovery",
            "Cashless payment adoption continued to strengthen after the initial shock",
        ],
        "Our overall conclusion is that COVID changed multiple dimensions of yellow taxi behavior, but those dimensions did not recover in the same way. Tipping rose during the disruption and intermediate period but moved back close to baseline by 2023, which makes it look more temporary. Airport recovery remained incomplete and uneven, with JFK recovering more strongly than LaGuardia in both the descriptive and regression evidence. Payment behavior shows the most persistent shift, with cashless usage continuing to grow well after the initial shock. So the main message we want the audience to remember is that recovery in the taxi system was not uniform: some behaviors reverted, some diverged, and some appear to have structurally shifted.",
        {"title": "Takeaway", "lines": ["Different behaviors recovered differently", "The recovery was not uniform across the taxi system"], "fill": LIGHT_GOLD},
    )

    add_text_slide(
        prs,
        13,
        "Backup: RQ2 Coefficient Table",
        "Backup",
        "Backup",
        [
            "JFK x COVID: +0.0005, p = 0.9175",
            "JFK x intermediate: +0.0108, p = 0.0031",
            "JFK x post-COVID: +0.0125, p < 1e-9",
            "Interpretation: later recovery terms are positive and significant, supporting stronger JFK recovery relative to LaGuardia",
        ],
        "Use this slide if the audience asks about the exact inferential evidence behind the airport claim.",
        {"title": "Model reminder", "lines": ["Outcome: airport trip share", "Reference categories: LaGuardia and 2019", "Month fixed effects included"], "fill": LIGHT_BLUE},
    )

    add_backup_metrics_slide(prs, 14)

    add_text_slide(
        prs,
        15,
        "Backup: Replication Workflow",
        "Backup",
        "Backup",
        [
            "The repository uses a shared DuckDB query layer and separate RQ folders",
            "Main commands:",
            (1, "python3 RQ1/tipping_analysis.py"),
            (1, "python3 RQ2/airport_trip_analysis.py"),
            (1, "python3 RQ3/cashless_payment_analysis.py"),
            "Local data requirement: monthly TLC parquet files plus `taxi_zone_lookup.csv` in `taxi_data/`",
        ],
        "Use this slide if someone asks how the work can be rerun from the repository.",
        {"title": "Submission reminder", "lines": ["Export slide deck as PDF", "Submit the PDF and video", "Keep the recording under 15 minutes"], "fill": LIGHT_GOLD},
    )

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_deck()
    print(f"Saved presentation draft to {OUTPUT_PATH}")
